import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def add_title(slide, text, left, top, width, height, font_size=18, color=RGBColor(0, 0, 0), bold=True, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tb

def draw_arrow(slide, start_x, start_y, end_x, end_y, direction='right', text="", color=RGBColor(16, 185, 129)):
    if direction == 'right':
        shape_type = MSO_SHAPE.RIGHT_ARROW
        arrow_width = end_x - start_x if end_x > start_x else Inches(0.5)
        arrow_height = Inches(0.15)
    elif direction == 'left':
        shape_type = MSO_SHAPE.LEFT_ARROW
        arrow_width = start_x - end_x if start_x > end_x else Inches(0.5)
        arrow_height = Inches(0.15)
        # For left arrow, the shape's left bounding box is end_x
        start_x = end_x
    else:
        shape_type = MSO_SHAPE.DOWN_ARROW
        arrow_width = Inches(0.15)
        arrow_height = end_y - start_y if end_y > start_y else Inches(0.5)
        
    arrow = slide.shapes.add_shape(
        shape_type, start_x, start_y, arrow_width, arrow_height
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    
    if text:
        if direction == 'right':
            mid_x = start_x + arrow_width / 2
            mid_y = start_y - Inches(0.3)
            add_title(slide, text, mid_x - Inches(1.5), mid_y, Inches(3.0), Inches(0.4), font_size=12, color=color, align=PP_ALIGN.CENTER)
        elif direction == 'left':
            mid_x = start_x + arrow_width / 2
            mid_y = start_y - Inches(0.3)
            add_title(slide, text, mid_x - Inches(1.5), mid_y, Inches(3.0), Inches(0.4), font_size=12, color=color, align=PP_ALIGN.CENTER)
        else:
            mid_x = start_x
            mid_y = start_y + arrow_height / 2
            add_title(slide, text, mid_x + Inches(0.1), mid_y - Inches(0.2), Inches(2.0), Inches(0.4), font_size=12, color=color, align=PP_ALIGN.LEFT)

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()

    add_title(slide, "LPO PAM4 System Optimization: Architecture & Evaluation", 
              Inches(0.5), Inches(0.1), Inches(12.333), Inches(0.5), 
              font_size=22, color=RGBColor(15, 23, 42), align=PP_ALIGN.CENTER)

    dsp_img = "scratch/dsp_arch.png"
    opt_img = "scratch/opt_arch.png"
    # FIXED: using the 4-algorithm comparison image from 0716
    chart_img = "result/20260716_174952_comparison/mlse_ber_convergence.png"
    tx_eye = "diagnostic_results/112G/01b_After_Tx_CTLE.png"
    rx_eye = "diagnostic_results/112G/10_ADC_Out_2sps.png"

    # Optimization Engine (Top Left)
    if os.path.exists(opt_img):
        slide.shapes.add_picture(opt_img, Inches(0.5), Inches(0.7), width=Inches(5.5))
        add_title(slide, "1. Optimization Loop Architecture", Inches(0.5), Inches(0.5), Inches(5.5), Inches(0.3), font_size=13, color=RGBColor(30, 41, 59))
        
    # Results Chart (Top Right - 4 algorithms!)
    if os.path.exists(chart_img):
        slide.shapes.add_picture(chart_img, Inches(7.3), Inches(0.7), width=Inches(5.5))
        add_title(slide, "3. Global Convergence Comparison (4 Algos)", Inches(7.3), Inches(0.5), Inches(5.5), Inches(0.3), font_size=13, color=RGBColor(30, 41, 59))

    # DSP Link (Bottom Center)
    if os.path.exists(dsp_img):
        slide.shapes.add_picture(dsp_img, Inches(3.6), Inches(3.5), width=Inches(6.1))
        add_title(slide, "2. Physical Data Link (3-Column Layout)", Inches(3.6), Inches(3.2), Inches(6.1), Inches(0.3), font_size=13, color=RGBColor(30, 41, 59))

    # Tx Eye (Bottom Left, pointing to Tx Host)
    if os.path.exists(tx_eye):
        slide.shapes.add_picture(tx_eye, Inches(0.5), Inches(4.2), width=Inches(2.5))
        add_title(slide, "Tx Eye (CTLE Out)", Inches(0.5), Inches(3.9), Inches(2.5), Inches(0.3), font_size=12, color=RGBColor(71, 85, 105), align=PP_ALIGN.CENTER)
        # Arrow pointing RIGHT to the DSP diagram (Tx column)
        draw_arrow(slide, Inches(3.1), Inches(5.5), Inches(3.5), Inches(5.5), direction='right', color=RGBColor(217, 70, 235))

    # Rx Eye (Bottom Right, pointing to Rx Host)
    if os.path.exists(rx_eye):
        slide.shapes.add_picture(rx_eye, Inches(10.3), Inches(4.2), width=Inches(2.5))
        add_title(slide, "Rx Eye (ADC Out)", Inches(10.3), Inches(3.9), Inches(2.5), Inches(0.3), font_size=12, color=RGBColor(71, 85, 105), align=PP_ALIGN.CENTER)
        # Arrow pointing LEFT to the DSP diagram (Rx column)
        draw_arrow(slide, Inches(10.2), Inches(4.5), Inches(9.8), Inches(4.5), direction='left', color=RGBColor(217, 70, 235))

    # Architecture Connection Arrows
    # From Opt to DSP
    draw_arrow(slide, Inches(3.2), Inches(2.5), Inches(3.2), Inches(3.1), direction='down', text="Evaluate Params", color=RGBColor(16, 185, 129))
    # From DSP to Results
    draw_arrow(slide, Inches(7.0), Inches(2.5), Inches(7.0), Inches(3.1), direction='down', text="Feedback BER", color=RGBColor(249, 115, 22))

    prs.save("System_Optimization_Report_v4.pptx")
    print("Successfully generated System_Optimization_Report_v4.pptx")

if __name__ == '__main__':
    main()
